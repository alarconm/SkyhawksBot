import os
import nltk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
def parse_generated_text(generated_text):
    # Remove the square brackets from the generated text
    generated_text = generated_text.strip("[]")
    
    # Split the generated text into slide information
    slides = generated_text.split(", ")

    # Extract slide titles and content
    slide_titles = []
    for slide in slides:
        slide_title = slide.split(": ")[1]
        slide_titles.append(slide_title)

    return slide_titles

def generate_pptx(generated_text):
    # Parse the generated text into slide titles and content
    slide_titles = parse_generated_text(generated_text)

    # Create a PowerPoint presentation
    pptx = Presentation()

    # Add a slide for each title and content
    for slide_title in slide_titles:
        slide_layout = pptx.slide_layouts[5]  # Use a layout with a title and content placeholder
        slide = pptx.slides.add_slide(slide_layout)

        # Add the slide title
        title = slide.shapes.title
        title.text = slide_title

        # Customize the title text size and font
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.bold = True

    # Save the PowerPoint presentation
    filename = "generated_presentation.pptx"
    pptx.save(filename)

    return filename
